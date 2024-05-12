import os
from pptx import Presentation

def change_title_in_ppt(folder_name, new_title):
    folder_path = rf"C:\Users\ox060\Downloads\【國中資訊科技】姿勢辨識\【國中資訊科技】姿勢辨識\{folder_name}"
    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            ppt_path = os.path.join(folder_path, filename)
            prs = Presentation(ppt_path)
            first_slide = prs.slides[0]
            first_slide.shapes.title.text = new_title
            prs.save(ppt_path)

# 輸入(資料夾名稱, 修改的標題)的列表
folders_and_titles = [
    ("1-1 AI 姿勢辨識與生活上的應用", "1-1 AI 姿勢辨識與生活上的應用"),
    ("2-1 什麼是PoseNet", "2-1 什麼是PoseNet"),
    ("2-2 使用「PoseNet2Scratch」擴充功能", "2-2 使用「PoseNet2Scratch」擴充功能"),
    ("2-3 繪製人體骨架", "2-3 繪製人體骨架"),
    ("3-1 系統拆解", "3-1 系統拆解"),
    ("3-2 使用「PoseNet2Scratch」擴充功能", "3-2 使用「PoseNet2Scratch」擴充功能"),
    ("3-3 設定初始值", "3-3 設定初始值"),
    ("3-4 遊戲說明", "3-4 遊戲說明"),
    ("3-5 開始計時", "3-5 開始計時"),
    ("3-6 公佈分數", "3-6 公佈分數"),
    ("4-1 認識 Teachable Machine", "4-1 認識 Teachable Machine"),
    ("4-2 開啟建立模型介面", "4-2 開啟建立模型介面"),
    ("4-3 輸入範例 - 網路攝影機", "4-3 輸入範例 - 網路攝影機"),
    ("4-4 輸入範例 - 影像檔案", "4-4 輸入範例 - 影像檔案"),
    ("4-5 進行訓練", "4-5 進行訓練"),
    ("4-6 檢視與輸出模型", "4-6 檢視與輸出模型"),
    ("5-1 系統拆解", "5-1 系統拆解"),
    ("5-2 連接姿勢分類模型", "5-2 連接姿勢分類模型"),
    ("5-3 連接Google試算表", "5-3 連接Google試算表"),
    ("5-4 製作出題按鈕廣播", "5-4 製作出題按鈕廣播"),
    ("5-5 製作答題按鈕廣播", "5-5 製作答題按鈕廣播")
]

for folder, new_title in folders_and_titles:
    change_title_in_ppt(folder, new_title)
